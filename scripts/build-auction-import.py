#!/usr/bin/env python3
"""
Build the auction player import CSV from the source PPT (authority) joined to
the 7-column sheet (base price / granular role / prior IPL team only).

Why the PPT is authority: the Codex-merged CSV dropped players entirely,
dropped stats for players who have them on their slide, and invented 5
duplicate rows with conflicting base prices. Each of those players appears
exactly once in the deck, so rebuilding from slides resolves all of it.
See .claude/plans/binary-pondering-starfish.md for the full findings.

Read-only against both sources; writes one CSV plus a validation report to
stdout. Nothing here talks to the database — import happens through the
existing admin UI so it goes down the already-tested code path.

Output headers match IMPORT_COLUMN_ALIASES in src/lib/validation/auction.ts
exactly. The first seven map to real `players` columns; every later column
falls through parseImportRow into `stats` and auto-creates a
player_stat_definitions row.

STAT LAYOUTS — the deck uses three, and python-pptx returns label shapes in
z-order, NOT visual order, so the labels cannot be zipped positionally with
the values. The value blob's own order is what's fixed. Verified against
known IPL figures (Jadeja, Cummins, Rashid Khan, Bairstow, Shanaka):

  all-rounder (labels have BOTH 'wickets taken' AND runs) — 6 slots:
      matches, wickets, runs, economy, best bowling, batting average
      e.g. Jadeja  '268 160 3525 7.71 5/16 29.38'
  bowler (labels have 'wickets taken', no runs) — 4 slots:
      matches, wickets, economy, best bowling
      e.g. Sumit Kumar  '4 0 11.40 N/A'
  batter (labels have runs/average, no wickets) — 4 slots:
      matches, runs, batting average, then strike rate OR highest score
      depending on which label the slide carries
      e.g. Shai Hope  '9 183 22.8 150'

'-' and 'N/A' are genuine "not applicable" placeholders and must be treated
as value tokens, otherwise they leak into the extracted player name.
"""

import csv
import difflib
import re
import sys
from pptx import Presentation

PPT = "Copy of AUCTION PPT .pptx"
SHEET = "merged_auction_players_updated.csv"
OUT = "auction_players_import.csv"

# Players present in the deck but absent from the 7-column sheet have no base
# price; they are still auctionable lots, so they import at the lowest tier
# and are listed loudly in the report for the admin to price by hand.
DEFAULT_BASE_PRICE = 2_000_000  # 20 L

# --------------------------------------------------------------------------
# Normalisation tables
# --------------------------------------------------------------------------

# is_overseas is computed in SQL as `lower(nationality) <> 'india'`
# (20260730080000_auction.sql), so an unnormalised "INDIAN" would silently
# mark Indian players overseas and eat into every squad's overseas cap.
COUNTRIES = [
    "INDIA", "AUSTRALIA", "ENGLAND", "SOUTH AFRICA", "NEW ZEALAND", "WEST INDIES",
    "SRI LANKA", "AFGHANISTAN", "BANGLADESH", "ZIMBABWE", "IRELAND", "USA",
    "NETHERLANDS", "NEPAL", "SCOTLAND", "NAMIBIA", "OMAN", "UAE", "CANADA",
]
COUNTRY_FIX = {
    "INDIAN": "INDIA", "AUSTRAILIA": "AUSTRALIA", "AFGHANISTHAN": "AFGHANISTAN",
    "AFGANISTAN": "AFGHANISTAN", "AFGHATISTHAN": "AFGHANISTAN", "NEWZEALAND": "NEW ZEALAND",
    "NEW ZELAND": "NEW ZEALAND", "NEW ZEALND": "NEW ZEALAND", "SRILANKA": "SRI LANKA",
    "AMERICA": "USA", "ENGALND": "ENGLAND", "SOUTHAFRICA": "SOUTH AFRICA",
    "WESTINDIES": "WEST INDIES",
}

IPL_TEAM_FIX = {
    "KKR": "KOLKATA KNIGHT RIDERS",
    "RCB": "ROYAL CHALLENGERS BENGALURU",
    "ROYAL CHALLENGERS BANGALORE": "ROYAL CHALLENGERS BENGALURU",
    "ROYAL CHALLENGERS BANAGLORE": "ROYAL CHALLENGERS BENGALURU",
    "SRH": "SUNRISERS HYDERABAD",
    "GUJRAT TITANS": "GUJARAT TITANS",
    "LUCNOW SUPER GIANTS": "LUCKNOW SUPER GIANTS",
    # "not played yet" is an absence, not a team — blank so nothing renders it
    # as a franchise.
    "NPY": "",
    "NOT PLAYED YET": "",
}

ROLE_ORDER = ["WICKET KEEPER", "ALL ROUNDER", "BOWLER", "BATTER"]


def canonical_role(raw: str) -> str:
    """Collapse ~29 deck spellings to the 4 values role_limits can cap on."""
    r = re.sub(r"\s+", " ", raw.upper()).strip()
    if "KEEPER" in r:
        return "WICKET KEEPER"
    if "ALL" in r and "ROUND" in r:
        return "ALL ROUNDER"
    if any(k in r for k in ("BOWL", "SPIN", "PACE", "FAST", "CHINAMAN", "MEDIUM", "OFF BREAK")):
        return "BOWLER"
    if "BAT" in r:
        return "BATTER"
    return ""


def parse_price(raw: str):
    """'2 CR' -> 20000000, '75l' -> 7500000. Rupees, matching Money's INR render."""
    s = re.sub(r"\s+", "", raw.upper())
    m = re.match(r"^([\d.]+)(CR|CRORE|CRORES|L|LAKH|LAKHS)$", s)
    if not m:
        return None
    n = float(m.group(1))
    return int(n * 10_000_000) if m.group(2).startswith("CR") else int(n * 100_000)


def norm_name(n: str) -> str:
    return re.sub(r"[^A-Z ]", "", re.sub(r"\s+", " ", n.upper())).strip()


def canon_country(raw: str) -> str:
    c = norm_name(raw)
    c = COUNTRY_FIX.get(c, c)
    if c in COUNTRIES:
        return c
    close = difflib.get_close_matches(c, COUNTRIES, n=1, cutoff=0.80)
    return close[0] if close else c


def looks_like_country(text: str) -> bool:
    c = norm_name(text)
    if not c or len(c.split()) > 3:
        return False
    if c in COUNTRIES or c in COUNTRY_FIX:
        return True
    return bool(difflib.get_close_matches(c, COUNTRIES, n=1, cutoff=0.80))


# --------------------------------------------------------------------------
# PPT extraction
# --------------------------------------------------------------------------

LABEL_RE = re.compile(
    r"player career|total|average|economy|best bowling|highest|scored|taken|played|"
    r"matches|innings|wickets|runs|ipl auction|pot\s*\d|strike|rate|career",
    re.I,
)
# '-' and 'N/A' are real placeholders in this deck; '*' marks a not-out score.
VALUE_TOKEN_RE = re.compile(r"^(?:\d+(?:\.\d+)?(?:/\d+)?\*?|-+|N/?A)$", re.I)


def walk(shapes):
    """Stat labels live inside grouped shapes — a flat walk finds only 183 of 249."""
    for sh in shapes:
        sub = getattr(sh, "shapes", None)
        if sub is not None:
            try:
                yield from walk(sub)
                continue
            except Exception:
                pass
        yield sh


def slide_texts(slide):
    out = []
    for sh in walk(slide.shapes):
        if getattr(sh, "has_text_frame", False):
            t = sh.text_frame.text.strip()
            if t:
                # One slide writes best-bowling as '3 /33'; without this the
                # shape stops looking like pure values and leaks into the name.
                out.append(re.sub(r"\s*/\s*", "/", t))
    return out


def is_value_shape(text: str) -> bool:
    parts = text.split()
    if not parts:
        return False
    return all(VALUE_TOKEN_RE.match(p) for p in parts)


def value_tokens(texts):
    toks = []
    for t in texts:
        if is_value_shape(t):
            toks.extend(t.split())
    return toks


def clean_stat(v: str):
    """'19.45*' -> '19.45'; '-' / 'N/A' -> None (omit the key entirely)."""
    if v is None:
        return None
    v = v.strip().rstrip("*")
    if not v or re.match(r"^-+$", v) or re.match(r"^N/?A$", v, re.I):
        return None
    return v


def extract_players():
    prs = Presentation(PPT)
    players, dividers, skipped = [], [], []
    pool, pool_seq = None, 0

    for idx, slide in enumerate(prs.slides, 1):
        texts = slide_texts(slide)
        if not texts:
            skipped.append((idx, "empty slide"))
            continue
        blob = "\n".join(texts)
        low = blob.lower()

        if "ipl auction" in low and "player career" not in low and len(texts) <= 4:
            name = " ".join(t for t in texts if t.strip().lower() != "ipl auction")
            name = re.sub(r"POT\s*\d+", "", name, flags=re.I)
            name = re.sub(r"\s+", " ", name.replace("\n", " ")).strip(" -·")
            if name:
                pool_seq += 1
                # The deck's own numbering repeats "POT 12"; renumber by slide
                # order so a lexical sort equals true auction order.
                pool = f"POT {pool_seq:02d} · {name.upper()}"
                dividers.append((idx, pool))
            continue

        toks = value_tokens(texts)
        if "player career" not in low and len(toks) < 3:
            skipped.append((idx, blob[:60].replace("\n", " ")))
            continue

        name_parts, country = [], ""
        for t in texts:
            if LABEL_RE.search(t) or is_value_shape(t):
                continue
            if looks_like_country(t):
                country = country or canon_country(t)
                continue
            name_parts.append(re.sub(r"\s+", " ", t.replace("\n", " ")).strip())

        # Some slides carry the same name shape twice ('SAM CURRAN SAM CURRAN').
        seen, uniq = set(), []
        for part in name_parts:
            k = norm_name(part)
            if k and k not in seen:
                seen.add(k)
                uniq.append(part)
        name = re.sub(r"\s+", " ", " ".join(uniq)).strip()

        players.append({
            "slide": idx, "pool": pool or "", "name": name, "country": country,
            "tokens": toks,
            "has_wickets": "wickets taken" in low,
            "has_runs": bool(re.search(r"runs?\s+scored?|total\s+runs", low)),
            "has_sr": "strike rate" in low,
            "has_highest": "highest" in low,
            "has_avg": "average" in low,
        })

    return players, dividers, skipped


def batter_tail_keys(p, toks):
    """
    Name the batter layout's 3rd and 4th slots. A batting average carries a
    decimal and sits under ~80; a highest score is a whole number; a strike
    rate runs well above 100. Verified against slides 21, 215, 222, 247.
    """
    tail = [clean_stat(x) for x in toks[2:4]]
    a, b = (tail + [None, None])[:2]

    def val(x):
        try:
            return float(x)
        except (TypeError, ValueError):
            return None

    va, vb = val(a), val(b)

    # 'strike rate' + 'highest score' together: the larger value is the rate.
    if p["has_sr"] and p["has_highest"] and va is not None and vb is not None:
        return ["Strike Rate", "Highest Score"] if va >= vb else ["Highest Score", "Strike Rate"]
    if p["has_sr"] and not p["has_highest"]:
        return ["Batting Average", "Strike Rate"]

    # Otherwise it's average vs highest score — the decimal one is the average.
    dec_a = a is not None and "." in a
    dec_b = b is not None and "." in b
    if dec_a != dec_b:
        return ["Batting Average", "Highest Score"] if dec_a else ["Highest Score", "Batting Average"]
    if va is not None and vb is not None:
        return ["Highest Score", "Batting Average"] if va >= vb else ["Batting Average", "Highest Score"]
    return ["Batting Average", "Highest Score"]


def assign_stats(p, warnings):
    t = p["tokens"]
    label = p["name"] or f"slide {p['slide']}"
    stats, layout = {}, None

    if p["has_wickets"] and p["has_runs"]:
        layout = "allrounder"
        keys = ["Matches", "Wickets", "Runs", "Economy", "Best Bowling", "Batting Average"]
    elif p["has_wickets"]:
        layout = "bowler"
        keys = ["Matches", "Wickets", "Economy", "Best Bowling"]
    elif p["has_runs"] or p["has_sr"]:
        # The batter layout's last two slots are NOT in a fixed order across
        # the deck — 'highest score, average' and 'average, highest score'
        # both occur (slides 222 vs 21). Labels can't disambiguate (z-order),
        # so classify the two values by their own shape instead.
        layout = "batter"
        keys = ["Matches", "Runs"] + batter_tail_keys(p, t)
    else:
        # No labels survived the group walk — infer from token shape.
        if any("/" in x for x in t):
            layout, keys = "bowler(inferred)", ["Matches", "Wickets", "Economy", "Best Bowling"]
        elif len(t) >= 6:
            layout = "allrounder(inferred)"
            keys = ["Matches", "Wickets", "Runs", "Economy", "Best Bowling", "Batting Average"]
        elif len(t) >= 3:
            layout, keys = "batter(inferred)", ["Matches", "Runs", "Batting Average", "Highest Score"]
        else:
            warnings.append(f"{label}: could not infer stat layout from {t}")
            return {}, None

    if len(t) > len(keys):
        warnings.append(f"{label}: {len(t)} values for {layout} layout ({len(keys)} slots) {t}")

    for key, raw in zip(keys, t):
        v = clean_stat(raw)
        if v is None:
            continue
        if key == "Best Bowling":
            if "/" in v:
                w, _, r = v.partition("/")
                stats["Best Bowling Wickets"], stats["Best Bowling Runs"] = w, r
            continue
        stats[key] = v

    def fnum(k):
        try:
            return float(stats.get(k, ""))
        except (TypeError, ValueError):
            return None

    # Range sanity — catches a mis-assigned layout far better than eyeballing.
    for key, lo, hi in (("Matches", 0, 400), ("Economy", 3, 20),
                        ("Batting Average", 0, 80), ("Strike Rate", 40, 300)):
        v = fnum(key)
        if v is not None and not (lo <= v <= hi):
            warnings.append(f"{label}: {key}={v} outside {lo}-{hi} (layout={layout})")

    return stats, layout


# --------------------------------------------------------------------------
# Main
# --------------------------------------------------------------------------

def main():
    sheet = list(csv.DictReader(open(SHEET)))
    by_name = {}
    for r in sheet:
        by_name.setdefault(norm_name(r["Player Name"]), r)

    players, dividers, skipped = extract_players()
    warnings, rows, seen_ref, sheet_matched, unpriced, layouts = [], [], {}, set(), [], {}

    for p in players:
        key = norm_name(p["name"])
        row = by_name.get(key)
        match_note = None

        if row is None:
            # A stray fragment often survives ('TRAVIS HEAD ravis'), so a
            # sheet name contained in the slide name is a strong signal.
            subs = [k for k in by_name if k and (k in key or key in k) and len(k) > 6]
            if subs:
                key = max(subs, key=len)
                row, match_note = by_name[key], f"substring-matched -> {key!r}"
        if row is None:
            close = difflib.get_close_matches(key, list(by_name), n=1, cutoff=0.78)
            if close:
                key = close[0]
                row, match_note = by_name[key], f"fuzzy-matched -> {key!r}"

        if match_note:
            warnings.append(f"{p['name']!r} (slide {p['slide']}): {match_note}")

        if row is not None:
            sheet_matched.add(key)
            price = parse_price(row["Base Price"])
            if price is None:
                warnings.append(
                    f"{p['name']!r} (slide {p['slide']}): unparseable base price "
                    f"{row['Base Price']!r} — defaulted"
                )
                price = DEFAULT_BASE_PRICE
            role_raw = re.sub(r"\s+", " ", row["Role"].upper()).strip()
            ipl_raw = re.sub(r"\s+", " ", row["PY IPL Team"].upper()).strip()
            country = canon_country(row["Country"])
            if p["country"] and p["country"] != country:
                warnings.append(
                    f"{p['name']!r} (slide {p['slide']}): country differs — "
                    f"sheet={country} deck={p['country']} (using sheet)"
                )
        else:
            # In the deck but not the sheet: a real lot with no price yet.
            price, role_raw, ipl_raw = DEFAULT_BASE_PRICE, "", ""
            country = p["country"] or "INDIA"
            unpriced.append((p["slide"], p["name"], p["pool"]))

        role = canonical_role(role_raw)
        if not role:
            # Fall back to the pool name, which encodes the role group
            # ('CAPPED BOWLERS 2', 'UNCAPPED BATSMAN'), then to the stat shape.
            role = canonical_role(p["pool"])
        if not role:
            role = "BOWLER" if p["has_wickets"] and not p["has_runs"] else "ALL ROUNDER"
            warnings.append(f"{p['name']!r} (slide {p['slide']}): role inferred as {role}")

        ipl = IPL_TEAM_FIX.get(ipl_raw, ipl_raw)

        # external_ref makes admin_import_players idempotent — a corrected
        # re-upload updates instead of duplicating, which matters when there
        # is no time for a clean retry.
        ref = (key or norm_name(p["name"])).replace(" ", "-").lower()
        if ref in seen_ref:
            warnings.append(
                f"{p['name']!r} (slide {p['slide']}): duplicate of slide {seen_ref[ref]} — skipped"
            )
            continue
        seen_ref[ref] = p["slide"]

        stats, layout = assign_stats(p, warnings)
        layouts[layout] = layouts.get(layout, 0) + 1
        if not stats:
            warnings.append(f"{p['name']!r} (slide {p['slide']}): no stats extracted")

        rows.append({
            "External Ref": ref,
            "Player Name": p["name"].upper(),
            "Role": role,
            "Base Price": price,
            "Pool": p["pool"] or "POT 99 · UNASSIGNED",
            "Country": country,
            "IPL Team": ipl,
            "Playing Style": role_raw,
            **stats,
        })

    headers = [
        "External Ref", "Player Name", "Role", "Base Price", "Pool", "Country",
        "IPL Team", "Playing Style", "Matches", "Runs", "Batting Average",
        "Highest Score", "Strike Rate", "Wickets", "Economy",
        "Best Bowling Wickets", "Best Bowling Runs",
    ]
    with open(OUT, "w", newline="") as fh:
        w = csv.DictWriter(fh, fieldnames=headers, extrasaction="ignore")
        w.writeheader()
        w.writerows(rows)

    # ---------------- report ----------------
    print(f"PPT player slides : {len(players)}")
    print(f"Sheet rows        : {len(sheet)}")
    print(f"ROWS WRITTEN      : {len(rows)}  -> {OUT}")
    print(f"Stat layouts      : {layouts}")
    print()
    print("Pools in auction order:")
    for _, d in dividers:
        print(f"  {d}  ({sum(1 for r in rows if r['Pool'] == d)} players)")

    print("\nRole distribution:", {r: sum(1 for x in rows if x["Role"] == r) for r in ROLE_ORDER})
    overseas = sum(1 for r in rows if r["Country"] != "INDIA")
    print(f"Indian: {len(rows) - overseas} | Overseas: {overseas}")
    print("Countries:", sorted({r["Country"] for r in rows}))
    total = sum(r["Base Price"] for r in rows)
    print(f"Total base price: Rs {total:,} ({total / 10_000_000:.1f} Cr)")

    for k in ("Matches", "Runs", "Wickets", "Economy", "Batting Average",
              "Highest Score", "Strike Rate", "Best Bowling Wickets"):
        print(f"  {k:22s} populated: {sum(1 for r in rows if r.get(k))}/{len(rows)}")

    nostats = [r["Player Name"] for r in rows
               if not any(r.get(k) for k in ("Matches", "Runs", "Wickets", "Economy"))]
    print(f"\nPlayers with no stats at all: {len(nostats)} {nostats}")

    print(f"\nIN DECK BUT NOT IN SHEET — priced at 20 L, admin must review: {len(unpriced)}")
    for s, n, pool in unpriced:
        print(f"  slide {s:>3} [{pool}] {n!r}")

    orphans = [by_name[k]["Player Name"] for k in by_name if k not in sheet_matched]
    print(f"\nIN SHEET BUT NO DECK SLIDE (not imported): {len(orphans)}")
    for o in orphans:
        print(f"  {o!r}")

    print(f"\nWARNINGS: {len(warnings)}")
    for wn in warnings:
        print(f"  ! {wn}")


if __name__ == "__main__":
    sys.exit(main())
